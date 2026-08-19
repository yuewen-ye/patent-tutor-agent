"""Patent-course visual components built from native editable PowerPoint objects."""

from __future__ import annotations

from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from backend.app.presentation.renderer.canvas import Canvas
from backend.app.presentation.renderer.shapes import bullets, color, line, rect, text_box
from backend.app.presentation.renderer.theme import Theme


def legal_citation_card(slide, canvas: Canvas, theme: Theme, x: float, y: float, w: float, h: float, reference: str, summary: str, meaning: str | None = None) -> None:
    rect(slide, canvas, x, y, w, h, theme.surface)
    line(slide, canvas, x, y, 0.12, h, theme.accent, 0)
    text_box(slide, canvas, "法条锚定", x + 0.3, y + 0.25, w - 0.6, 0.3, theme=theme, size=12, bold=True, fill=theme.accent)
    text_box(slide, canvas, reference, x + 0.3, y + 0.72, w - 0.6, 0.48, theme=theme, size=22, bold=True, fill=theme.primary)
    text_box(slide, canvas, summary, x + 0.3, y + 1.35, w - 0.6, max(0.5, h - 2.05), theme=theme, size=16)
    if meaning:
        text_box(slide, canvas, f"适用意义：{meaning}", x + 0.3, y + h - 0.55, w - 0.6, 0.3, theme=theme, size=11, fill=theme.muted)


def irac_flow(slide, canvas: Canvas, theme: Theme, values: list[tuple[str, str]], x: float, y: float, w: float, h: float) -> None:
    count = len(values)
    gap = 0.25
    box_w = (w - gap * (count - 1)) / count
    for index, (label, value) in enumerate(values):
        left = x + index * (box_w + gap)
        rect(slide, canvas, left, y, box_w, h, theme.surface)
        text_box(slide, canvas, label, left + 0.14, y + 0.18, box_w - 0.28, 0.3, theme=theme, size=14, bold=True, fill=theme.accent, align=PP_ALIGN.CENTER)
        text_box(slide, canvas, value, left + 0.16, y + 0.72, box_w - 0.32, h - 0.9, theme=theme, size=14, align=PP_ALIGN.CENTER)
        if index + 1 < count:
            connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(left + box_w), Inches(y + h / 2), Inches(left + box_w + gap), Inches(y + h / 2))
            connector.line.color.rgb = color(theme.accent)
            connector.line.width = Inches(0.025)


def patent_timeline(slide, canvas: Canvas, theme: Theme, steps: list[str], x: float, y: float, w: float) -> None:
    count = min(5, len(steps))
    if not count:
        return
    gap = w / count
    line(slide, canvas, x + 0.25, y + 0.45, w - 0.5, 0.05, theme.secondary)
    for index, value in enumerate(steps[:count]):
        center = x + gap * index + gap / 2
        rect(slide, canvas, center - 0.2, y + 0.25, 0.4, 0.4, theme.accent)
        text_box(slide, canvas, str(index + 1), center - 0.2, y + 0.35, 0.4, 0.16, theme=theme, size=10, bold=True, fill="FFFFFF", align=PP_ALIGN.CENTER)
        text_box(slide, canvas, value, center - min(0.75, gap / 2 - 0.05), y + 0.85, min(1.5, gap - 0.1), 0.9, theme=theme, size=13, align=PP_ALIGN.CENTER)


def comparison_matrix(slide, canvas: Canvas, theme: Theme, left_title: str, right_title: str, rows: list[tuple[str, str]], x: float, y: float, w: float, h: float) -> None:
    columns = [left_title, "判断", right_title]
    widths = [w * 0.37, w * 0.22, w * 0.37]
    current_x = x
    for header, width in zip(columns, widths, strict=True):
        rect(slide, canvas, current_x, y, width, 0.55, theme.primary)
        text_box(slide, canvas, header, current_x + 0.08, y + 0.16, width - 0.16, 0.2, theme=theme, size=13, bold=True, fill="FFFFFF", align=PP_ALIGN.CENTER)
        current_x += width
    row_h = min(0.72, (h - 0.6) / max(1, len(rows)))
    for row_index, (left, right) in enumerate(rows):
        values = [left, "对比", right]
        current_x = x
        for cell, width in zip(values, widths, strict=True):
            rect(slide, canvas, current_x, y + 0.55 + row_index * row_h, width, row_h, theme.surface)
            text_box(slide, canvas, cell, current_x + 0.08, y + 0.68 + row_index * row_h, width - 0.16, row_h - 0.15, theme=theme, size=12, align=PP_ALIGN.CENTER)
            current_x += width


def exam_question_card(slide, canvas: Canvas, theme: Theme, question: str, checklist: list[str], warning: str | None, x: float, y: float, w: float, h: float) -> None:
    rect(slide, canvas, x, y, w, h, theme.surface)
    text_box(slide, canvas, "练习题", x + 0.25, y + 0.2, w - 0.5, 0.3, theme=theme, size=13, bold=True, fill=theme.accent)
    text_box(slide, canvas, question, x + 0.25, y + 0.68, w - 0.5, 0.8, theme=theme, size=18, bold=True)
    bullets(slide, canvas, checklist[:4], x + 0.22, y + 1.65, w - 0.44, h - 2.1, theme=theme, size=14)
    if warning:
        text_box(slide, canvas, f"易错点：{warning}", x + 0.25, y + h - 0.42, w - 0.5, 0.22, theme=theme, size=11, fill=theme.warning)
