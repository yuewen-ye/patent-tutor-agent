"""Additive visual effects for the presentation renderer.

Activated ONLY when a slide's ``visual_intent`` free text carries an explicit
directive (``gradient:`` / ``illustration:``). When no directive is present, the
existing render paths are byte-for-byte unchanged.

No contract field, no existing function, and no database surface is modified.

Gradient fills are emulated by stacking thin solid-colour bands with RGB-linear
interpolation, so they work on any python-pptx version without relying on the
incomplete high-level gradient API. Flat-style illustrations are composed from
native auto-shapes (oval / triangle / rounded rectangle / pentagon / star /
connector), so no external image asset is required.
"""

from __future__ import annotations

import math
import re

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from backend.app.presentation.renderer.canvas import Canvas
from backend.app.presentation.renderer.shapes import color, line, rect, text_box
from backend.app.presentation.renderer.theme import Theme

_GRADIENT_RE = re.compile(
    r"gradient\s*:\s*([hvd])\s*(?:\(\s*([0-9A-Fa-f]{6})\s*(?:->|→|—|-)\s*([0-9A-Fa-f]{6})\s*\))?",
    re.IGNORECASE,
)
_ILLUSTRATION_RE = re.compile(r"illustration\s*:\s*([a-z_]+)", re.IGNORECASE)

# Templates that trigger the full "gradient hero cover" branch (early return).
_COVER_TEMPLATES = {"title", "cover_minimal", "cover_split", "hero_statement"}


def parse_visual_intent(intent: str) -> dict[str, object]:
    """Extract additive effect directives from a slide's ``visual_intent`` text."""
    if not intent:
        return {}
    out: dict[str, object] = {}
    g = _GRADIENT_RE.search(intent)
    if g:
        out["gradient"] = {
            "axis": g.group(1).lower(),
            "from": (g.group(2) or "").upper() or None,
            "to": (g.group(3) or "").upper() or None,
        }
    m = _ILLUSTRATION_RE.search(intent)
    if m:
        out["illustration"] = m.group(1).lower()
    return out


# --------------------------------------------------------------------------- #
# Banded gradient (uses only the existing solid-fill ``rect`` primitive)
# --------------------------------------------------------------------------- #
def _lerp(a: int, b: int, t: float) -> int:
    return max(0, min(255, round(a + (b - a) * t)))


def lerp_color(c1: str, c2: str, t: float) -> str:
    r1, g1, b1 = int(c1[0:2], 16), int(c1[2:4], 16), int(c1[4:6], 16)
    r2, g2, b2 = int(c2[0:2], 16), int(c2[2:4], 16), int(c2[4:6], 16)
    return f"{_lerp(r1, r2, t):02X}{_lerp(g1, g2, t):02X}{_lerp(b1, b2, t):02X}"


def draw_gradient(
    slide,
    canvas: Canvas,
    x: float,
    y: float,
    w: float,
    h: float,
    c_from: str,
    c_to: str,
    axis: str = "h",
    bands: int = 28,
) -> None:
    """Draw a banded gradient using only the existing solid-fill ``rect`` primitive.

    ``axis`` controls band direction: ``h`` = left→right, ``v`` = top→bottom,
    ``d`` = diagonal (emulated as horizontal). Adjacent bands overlap by 0.01 inch
    to avoid hairline seams from shape anti-aliasing.
    """
    n = max(2, bands)
    if axis == "v":
        bh = h / n
        for k in range(n):
            c = lerp_color(c_from, c_to, (k + 0.5) / n)
            rect(slide, canvas, x, y + k * bh, w, bh + 0.012, c, radius=False)
    else:  # 'h' or 'd' (diagonal emulated as horizontal for band simplicity)
        bw = w / n
        for k in range(n):
            c = lerp_color(c_from, c_to, (k + 0.5) / n)
            rect(slide, canvas, x + k * bw, y, bw + 0.012, h, c, radius=False)


def apply_gradient_background(slide, canvas: Canvas, theme: Theme, spec: dict) -> None:
    """Full-bleed gradient backdrop for non-cover templates."""
    axis = spec.get("axis", "h")
    c_from = spec.get("from") or theme.background
    c_to = spec.get("to") or theme.grid
    draw_gradient(
        slide, canvas, 0, 0, canvas.width, canvas.height, c_from, c_to, axis=axis, bands=32
    )


# --------------------------------------------------------------------------- #
# Flat-style vector illustrations composed from native auto-shapes
# --------------------------------------------------------------------------- #
def _prim(slide, canvas: Canvas, x, y, w, h, shape_enum, fill):
    shape = slide.shapes.add_shape(shape_enum, *canvas.box(x, y, w, h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.fill.background()
    return shape


def _oval(slide, canvas, x, y, w, h, fill):
    return _prim(slide, canvas, x, y, w, h, MSO_SHAPE.OVAL, fill)


def _tri(slide, canvas, x, y, w, h, fill):
    return _prim(slide, canvas, x, y, w, h, MSO_SHAPE.ISOCELES_TRIANGLE, fill)


def _rrect(slide, canvas, x, y, w, h, fill):
    return _prim(slide, canvas, x, y, w, h, MSO_SHAPE.ROUNDED_RECTANGLE, fill)


def _star(slide, canvas, x, y, w, h, fill):
    return _prim(slide, canvas, x, y, w, h, MSO_SHAPE.STAR_5_POINT, fill)


def _conn(slide, x1, y1, x2, y2, stroke, width_pt: float = 1.75):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.line.color.rgb = color(stroke)
    c.line.width = Pt(width_pt)
    return c


def _decor_motif(slide, canvas, theme: Theme, x, y, w, h) -> None:
    """Fallback emblem (ring + star + dots) when illustration kind is unknown."""
    cx, cy = x + w / 2, y + h / 2
    _oval(slide, canvas, cx - 1.25, cy - 1.25, 2.5, 2.5, theme.grid)
    _oval(slide, canvas, cx - 0.95, cy - 0.95, 1.9, 1.9, theme.surface)
    _star(slide, canvas, cx - 0.55, cy - 0.55, 1.1, 1.1, theme.accent)
    for dx, dy in [(-1.55, -1.05), (1.55, -1.05), (0, 1.45)]:
        _oval(slide, canvas, cx + dx - 0.12, cy + dy - 0.12, 0.24, 0.24, theme.accent)


def _ill_lightbulb(slide, canvas, theme, x, y, w, h):
    cx, cy = x + w / 2, y + h / 2
    d = min(w, h) * 0.42
    _oval(slide, canvas, cx - d / 2, cy - d / 2 - 0.2, d, d, theme.accent)
    _rrect(slide, canvas, cx - d * 0.18, cy + d / 2 - 0.25, d * 0.36, 0.5, theme.muted)
    _rrect(slide, canvas, cx - d * 0.28, cy + d / 2 + 0.3, d * 0.56, 0.22, theme.muted)
    for dx, dy in [(-d * 0.95, -d * 0.8), (d * 0.95, -d * 0.8), (0, -d * 1.25)]:
        _oval(slide, canvas, cx + dx - 0.12, cy + dy - 0.12, 0.24, 0.24, theme.success)


def _ill_scales(slide, canvas, theme, x, y, w, h):
    cx, cy = x + w / 2, y + h / 2
    span = w * 0.7
    _rrect(slide, canvas, cx - 0.6, cy + 1.2, 1.2, 0.22, theme.muted)
    line(slide, canvas, cx - 0.03, cy - 1.0, 0.06, 2.2, theme.muted)
    line(slide, canvas, cx - span / 2, cy - 0.9, span, 0.06, theme.accent)
    _oval(slide, canvas, cx - 0.12, cy - 1.12, 0.24, 0.24, theme.accent)
    _conn(slide, cx - span / 2, cy - 0.87, cx - span / 2, cy - 0.5, theme.muted)
    _oval(slide, canvas, cx - span / 2 - 0.3, cy - 0.5, 0.6, 0.3, theme.grid)
    _conn(slide, cx + span / 2, cy - 0.87, cx + span / 2, cy - 0.5, theme.muted)
    _oval(slide, canvas, cx + span / 2 - 0.3, cy - 0.5, 0.6, 0.3, theme.grid)


def _ill_path(slide, canvas, theme, x, y, w, h):
    pts = [
        (x + 0.25, y + h - 0.55),
        (x + w * 0.4, y + h * 0.55),
        (x + w * 0.72, y + h * 0.32),
        (x + w - 0.35, y + 0.45),
    ]
    for i in range(len(pts) - 1):
        _conn(slide, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], theme.muted)
    for i, (px, py) in enumerate(pts):
        col = theme.accent if i == len(pts) - 1 else theme.grid
        _oval(slide, canvas, px - 0.2, py - 0.2, 0.4, 0.4, col)
        text_box(
            slide, canvas, str(i + 1), px - 0.2, py - 0.14, 0.4, 0.2,
            theme=theme, size=10, fill="FFFFFF", align=PP_ALIGN.CENTER,
        )


def _ill_document(slide, canvas, theme, x, y, w, h):
    cx = x + w / 2
    dw, dh = w * 0.55, h * 0.7
    dx, dy = cx - dw / 2, y + h * 0.15
    _rrect(slide, canvas, dx, dy, dw, dh, theme.surface)
    _tri(slide, canvas, dx + dw - 0.5, dy, 0.5, 0.5, theme.grid)
    line(slide, canvas, dx + 0.25, dy + 0.18, dw * 0.4, 0.08, theme.accent)
    for i in range(4):
        ly = dy + 0.3 + i * 0.42
        line(slide, canvas, dx + 0.25, ly, dw - 0.55, 0.05, theme.muted)


def _ill_book(slide, canvas, theme, x, y, w, h):
    cx, cy = x + w / 2, y + h / 2
    bw, bh = w * 0.32, h * 0.55
    _rrect(slide, canvas, cx - bw - 0.05, cy - bh / 2, bw, bh, theme.accent)
    _rrect(slide, canvas, cx + 0.05, cy - bh / 2, bw, bh, theme.grid)
    line(slide, canvas, cx - 0.03, cy - bh / 2, 0.06, bh, theme.muted)
    for i in range(3):
        ly = cy - bh / 2 + 0.22 + i * 0.32
        line(slide, canvas, cx - bw + 0.06, ly, bw - 0.22, 0.04, theme.surface)
        line(slide, canvas, cx + 0.16, ly, bw - 0.22, 0.04, theme.muted)


def _ill_concept(slide, canvas, theme, x, y, w, h):
    cx, cy = x + w / 2, y + h / 2
    _oval(slide, canvas, cx - 0.55, cy - 0.55, 1.1, 1.1, theme.accent)
    text_box(
        slide, canvas, "核心", cx - 0.55, cy - 0.22, 1.1, 0.44,
        theme=theme, size=12, fill="FFFFFF", align=PP_ALIGN.CENTER,
    )
    r = min(w, h) * 0.38
    for a in (0.0, math.pi / 2, math.pi, 3 * math.pi / 2):
        sx, sy = cx + math.cos(a) * r, cy + math.sin(a) * r
        _conn(slide, cx, cy, sx, sy, theme.muted)
        _oval(slide, canvas, sx - 0.3, sy - 0.3, 0.6, 0.6, theme.grid)


def _ill_star(slide, canvas, theme, x, y, w, h):
    cx, cy = x + w / 2, y + h / 2
    d = min(w, h) * 0.55
    _star(slide, canvas, cx - d / 2, cy - d / 2, d, d, theme.accent)
    d2 = d * 0.6
    _star(slide, canvas, cx - d2 / 2, cy - d2 / 2 + 0.05, d2, d2, theme.grid)
    for dx, dy in [(-d * 0.85, d * 0.4), (d * 0.85, d * 0.4), (0, -d * 0.95)]:
        _oval(slide, canvas, cx + dx - 0.1, cy + dy - 0.1, 0.2, 0.2, theme.success)


_ILLUSTRATIONS = {
    "lightbulb": _ill_lightbulb, "idea": _ill_lightbulb,
    "scales": _ill_scales, "balance": _ill_scales,
    "path": _ill_path, "journey": _ill_path,
    "document": _ill_document, "filing": _ill_document,
    "book": _ill_book, "learning": _ill_book,
    "concept": _ill_concept, "hub": _ill_concept,
    "star": _ill_star, "achievement": _ill_star,
}


def draw_illustration(slide, canvas, theme: Theme, kind: str, x, y, w, h) -> None:
    fn = _ILLUSTRATIONS.get(kind)
    if fn:
        fn(slide, canvas, theme, x, y, w, h)
    else:
        _decor_motif(slide, canvas, theme, x, y, w, h)


# --------------------------------------------------------------------------- #
# Gradient hero cover (NEW full-render branch; existing ``cover()`` untouched)
# --------------------------------------------------------------------------- #
def _gradient_hero_cover(slide, canvas, item, theme, page, effects) -> None:
    grad = effects.get("gradient")
    spec = grad if isinstance(grad, dict) else {}
    axis = spec.get("axis", "h")
    c_from = spec.get("from") or theme.background
    c_to = spec.get("to") or theme.grid
    draw_gradient(
        slide, canvas, 0, 0, canvas.width, canvas.height, c_from, c_to, axis=axis, bands=32
    )
    # Left title card (white rounded) — mirrors brand: cream bg + white card + orange rule
    rect(slide, canvas, 0.72, 1.05, 7.5, 4.85, theme.surface)
    line(slide, canvas, 1.0, 1.42, 2.6, 0.06, theme.accent)
    text_box(
        slide, canvas, item.title, 1.0, 1.65, 6.9, 2.2,
        theme=theme, size=36, bold=False, fill=theme.primary,
    )
    if item.subtitle:
        text_box(
            slide, canvas, item.subtitle, 1.0, 3.95, 6.9, 0.9,
            theme=theme, size=15, fill=theme.muted,
        )
    if item.legal_reference:
        text_box(
            slide, canvas, item.legal_reference, 1.0, 5.4, 6.9, 0.35,
            theme=theme, size=11, fill=theme.muted,
        )
    # Right illustration / motif zone
    illu = effects.get("illustration")
    if illu:
        draw_illustration(slide, canvas, theme, illu, 8.6, 1.5, 4.1, 4.0)
    else:
        _decor_motif(slide, canvas, theme, 8.6, 1.5, 4.1, 4.0)
    text_box(
        slide, canvas, f"PATENT TUTOR · {page:02d}", 1.0, 6.7, 6, 0.28,
        theme=theme, size=10, fill=theme.muted,
    )
    text_box(
        slide, canvas, f"{page:02d}", 12.2, 0.5, 0.7, 0.3,
        theme=theme, size=10, fill=theme.muted, align=PP_ALIGN.RIGHT,
    )


def maybe_render_visual_effects(
    slide, canvas, item, theme: Theme, page: int, template: str
) -> bool:
    """Additive dispatch hook.

    Returns True when the slide was fully rendered by an effect branch (caller
    should return immediately). Returns False when only a backdrop was drawn
    (caller should continue with the existing template handler) or when no
    effect directive was present. Existing paths are untouched when no
    directive is present.
    """
    effects = parse_visual_intent(item.visual_intent or "")
    grad = effects.get("gradient")
    illu = effects.get("illustration")
    if not grad and not illu:
        return False
    if template in _COVER_TEMPLATES:
        _gradient_hero_cover(slide, canvas, item, theme, page, effects)
        return True
    if isinstance(grad, dict):
        apply_gradient_background(slide, canvas, theme, grad)
    return False
