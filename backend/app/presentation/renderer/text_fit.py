"""Text fitting helpers to reduce overflow/overlap in generated slides.

PowerPoint files can declare ``TEXT_TO_FIT_SHAPE`` auto-fit, but headless
renderers (LibreOffice used for PNG previews) do not always honor it. We
therefore pre-calculate an approximate font scale so the rendered result is
close to what was intended.
"""

from __future__ import annotations

from pptx.util import Pt

# Approximate character widths (inches) for the fonts we use at 1 pt.
# CJK glyphs are roughly square; Latin glyphs average about half an em.
_CJK_CHAR_WIDTH_PT = 0.018
_LATIN_CHAR_WIDTH_PT = 0.009
_LINE_HEIGHT_FACTOR = 1.25
_MIN_FONT_SIZE = 9.0


def _estimated_char_width_pt(text: str) -> float:
    """Average glyph width for the given text in points."""
    if not text:
        return _LATIN_CHAR_WIDTH_PT
    cjk_count = sum(1 for ch in text if "\u4e00" <= ch <= "\u9fff")
    total = len(text)
    latin_count = total - cjk_count
    return (
        cjk_count * _CJK_CHAR_WIDTH_PT + latin_count * _LATIN_CHAR_WIDTH_PT
    ) / total


def estimate_text_height(
    text: str,
    width_inches: float,
    font_size_pt: float,
    line_height_factor: float = _LINE_HEIGHT_FACTOR,
) -> float:
    """Return an approximate rendered text height in inches.

    Word wrap is assumed. Empty text returns 0.
    """
    if not text or width_inches <= 0 or font_size_pt <= 0:
        return 0.0
    char_width = _estimated_char_width_pt(text) * font_size_pt
    chars_per_line = max(1, int(width_inches / char_width))
    lines = (len(text) + chars_per_line - 1) // chars_per_line
    line_height = font_size_pt * line_height_factor / 72.0
    return lines * line_height


def scaled_font_size_to_fit(
    text: str,
    width_inches: float,
    height_inches: float,
    requested_size_pt: float,
    min_size_pt: float = _MIN_FONT_SIZE,
) -> float:
    """Pick a font size that keeps the text inside the box.

    If the requested size already fits, it is returned unchanged. Otherwise the
    size is reduced (down to ``min_size_pt``) until the estimated height fits.
    """
    if not text or height_inches <= 0:
        return requested_size_pt
    size = requested_size_pt
    while size > min_size_pt:
        if estimate_text_height(text, width_inches, size) <= height_inches:
            return size
        size -= 1.0
    return min_size_pt


def apply_text_fit(
    text_frame,
    width_inches: float,
    height_inches: float,
    requested_size_pt: float,
    *,
    min_size_pt: float = _MIN_FONT_SIZE,
    auto_fit: bool = True,
) -> float:
    """Set font size on every run/paragraph in ``text_frame`` to a fitted size.

    Also sets the text frame's auto-fit property when ``auto_fit`` is true so
    that desktop PowerPoint can make a final adjustment if needed.

    Returns the font size actually applied (in points).
    """
    from pptx.enum.text import MSO_AUTO_SIZE

    text = text_frame.text or ""
    fitted = scaled_font_size_to_fit(
        text, width_inches, height_inches, requested_size_pt, min_size_pt
    )
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(fitted)
    if auto_fit:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return fitted
