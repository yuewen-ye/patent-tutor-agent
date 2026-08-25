"""Reusable native PowerPoint shape and text primitives."""

from __future__ import annotations

import re
from collections.abc import Iterable

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from backend.app.presentation.renderer.canvas import Canvas
from backend.app.presentation.renderer.text_fit import apply_text_fit, scaled_font_size_to_fit
from backend.app.presentation.renderer.theme import Theme


def color(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def rect(slide, canvas: Canvas, x: float, y: float, w: float, h: float, fill: str, radius: bool = True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, *canvas.box(x, y, w, h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.fill.background()
    return shape


def line(slide, canvas: Canvas, x: float, y: float, w: float, h: float, stroke: str, width: float = 1.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *canvas.box(x, y, w, h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(stroke)
    shape.line.fill.background()
    return shape


_CJK_RE = re.compile(r"[\u4e00-\u9fff\u3400-\u4dbf\u3040-\u30ff\u3000-\u303f\uff00-\uffef]")


def _estimate_text_units(text: str) -> float:
    """Approximate the text footprint: each CJK char = 1 unit, each ASCII word = 0.55 unit."""
    if not text:
        return 0.0
    cjk_count = len(_CJK_RE.findall(text))
    ascii_part = _CJK_RE.sub(" ", text)
    ascii_words = [w for w in ascii_part.split() if w]
    return float(cjk_count) + 0.55 * len(ascii_words)


def _auto_fit_font_size(size: float, text: str, w: float, h: float) -> float:
    """Shrink font size gradually when text overflows the text box."""
    if not text or size <= 9:
        return max(9.0, float(size))
    units = _estimate_text_units(text)
    if units <= 0:
        return max(9.0, float(size))
    chars_per_line = max(1.0, (w / 12.0) * 13.0 * (18.0 / size))
    lines_needed = units / chars_per_line
    line_height_in = 0.3 * (size / 18.0)
    height_needed_in = lines_needed * line_height_in
    if height_needed_in <= h and lines_needed <= max(2.0, h / line_height_in * 0.95):
        return max(9.0, float(size))
    ratio = max(height_needed_in / max(0.05, h), lines_needed / max(2.0, h / max(0.05, line_height_in)))
    new_size = size / (ratio ** 0.5)
    return max(9.0, min(size, new_size))


def _apply_cross_platform_font(font, theme: Theme) -> None:
    """Declare Latin and East-Asian typefaces in PPTX OOXML.

    ``python-pptx`` maps ``Font.name`` to the Latin typeface only. Chinese text
    therefore relied on the renderer's fallback font in headless LibreOffice.
    Set ``ea`` explicitly so LibreOffice can resolve the requested CJK font via
    fontconfig on Linux, while PowerPoint uses Microsoft YaHei on Windows.
    """
    font.name = theme.font
    font._element.set("ea", theme.cjk_font)


def text_box(
    slide,
    canvas: Canvas,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    theme: Theme,
    size: float = 18,
    bold: bool = False,
    fill: str | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
    fit_text: bool = True,
):
    safe_text = str(text or "").strip()
    inner_w = max(0.1, w - 2 * margin)
    inner_h = max(0.08, h - 2 * margin)
    fitted_size = _auto_fit_font_size(float(size), safe_text, inner_w, inner_h)
    box = slide.shapes.add_textbox(*canvas.box(x, y, w, h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = safe_text
    _apply_cross_platform_font(run.font, theme)
    run.font.size = Pt(fitted_size)
    run.font.bold = bold
    run.font.color.rgb = color(fill or theme.text)
    if fit_text:
        apply_text_fit(tf, w - 2 * margin, h - 2 * margin, size)
    return box


def bullets(
    slide,
    canvas: Canvas,
    items: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    theme: Theme,
    size: float = 18,
    fit_text: bool = True,
):
    box = slide.shapes.add_textbox(*canvas.box(x, y, w, h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    items = [str(it) for it in items if it is not None]
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.space_after = Pt(max(3, int(8 * (size / 18.0))))
        _apply_cross_platform_font(p.font, theme)
        p.font.size = Pt(size)
        p.font.color.rgb = color(theme.text)
    if fit_text and items:
        available_h = h - 0.08
        longest = max(items, key=len)
        fitted = max(
            9.0,
            min(
                size,
                size
                * (available_h / max(0.3, len(items) * size * 1.35 / 72.0)),
            ),
        )
        fitted = scaled_font_size_to_fit(longest, w - 0.13, available_h / max(1, len(items)), fitted)
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(fitted)
    return box
