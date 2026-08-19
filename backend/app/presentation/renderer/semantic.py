"""Semantic visual layer: meaning-driven diagrams, not generic bullet pages."""

from __future__ import annotations

from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from backend.app.presentation.contracts import PresentationVisualElement
from backend.app.presentation.renderer.canvas import Canvas
from backend.app.presentation.renderer.decor import add_callout
from backend.app.presentation.renderer.shapes import color, line, rect, text_box
from backend.app.presentation.renderer.theme import Theme


def render_element(slide, canvas: Canvas, theme: Theme, element: PresentationVisualElement, x: float, y: float, w: float, h: float) -> None:
    values = element.values or element.items
    if element.type == "callout":
        add_callout(slide, canvas, theme, element.title or element.label or "重点", element.text or (values[0] if values else ""), x, y, w, h)
    elif element.type == "warning_panel":
        add_callout(slide, canvas, theme, element.title or "易错点", element.text or (values[0] if values else ""), x, y, w, h, warning=True)
    elif element.type in {"timeline", "evidence_stack"}:
        count = max(1, min(6, len(values)))
        step_w = w / count
        line(slide, canvas, x + 0.2, y + 0.48, w - 0.4, 0.05, theme.secondary)
        for i, value in enumerate(values[:count]):
            cx = x + i * step_w + step_w / 2
            rect(slide, canvas, cx - 0.18, y + 0.28, 0.36, 0.36, theme.accent)
            text_box(slide, canvas, str(i + 1), cx - 0.18, y + 0.37, 0.36, 0.15, theme=theme, size=9, bold=True, fill="FFFFFF", align=PP_ALIGN.CENTER)
            text_box(slide, canvas, value, cx - step_w / 2 + 0.05, y + 0.85, step_w - 0.1, h - 0.9, theme=theme, size=12, align=PP_ALIGN.CENTER)
    elif element.type == "decision_tree":
        count = max(1, min(4, len(values)))
        node_w = min(2.5, w / count - 0.15)
        for i, value in enumerate(values[:count]):
            nx = x + i * (w / count) + 0.08
            rect(slide, canvas, nx, y + 0.35, node_w, 0.8, theme.surface)
            text_box(slide, canvas, value, nx + 0.12, y + 0.58, node_w - 0.24, 0.28, theme=theme, size=13, bold=True, align=PP_ALIGN.CENTER)
            if i:
                connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(nx - 0.12), Inches(y + 0.75), Inches(nx), Inches(y + 0.75))
                connector.line.color.rgb = color(theme.accent)
    elif element.type == "concept_map":
        center_x, center_y = x + w / 2, y + h / 2
        rect(slide, canvas, center_x - 1.0, center_y - 0.35, 2.0, 0.7, theme.primary)
        text_box(slide, canvas, element.title or "核心概念", center_x - 0.85, center_y - 0.12, 1.7, 0.22, theme=theme, size=14, bold=True, fill="FFFFFF", align=PP_ALIGN.CENTER)
        for i, value in enumerate(values[:4]):
            nx = x + (i % 2) * (w - 2.2) + 0.1
            ny = y + (i // 2) * (h - 1.1) + 0.1
            rect(slide, canvas, nx, ny, 2.1, 0.6, theme.surface)
            text_box(slide, canvas, value, nx + 0.1, ny + 0.2, 1.9, 0.2, theme=theme, size=12, align=PP_ALIGN.CENTER)
    elif element.type in {"comparison_matrix", "metric_cards"}:
        count = max(1, min(4, len(values)))
        card_w = (w - 0.2 * (count - 1)) / count
        for i, value in enumerate(values[:count]):
            cx = x + i * (card_w + 0.2)
            rect(slide, canvas, cx, y, card_w, h, theme.surface)
            text_box(slide, canvas, str(i + 1), cx + 0.15, y + 0.15, 0.35, 0.25, theme=theme, size=11, bold=True, fill=theme.accent)
            text_box(slide, canvas, value, cx + 0.15, y + 0.65, card_w - 0.3, h - 0.85, theme=theme, size=14, align=PP_ALIGN.CENTER)
    elif element.type == "irac":
        labels = ["Issue", "Rule", "Application", "Conclusion"]
        vals = values[:4] or [element.text or ""]
        for i, value in enumerate(vals):
            bx = x + i * (w / len(vals))
            rect(slide, canvas, bx + 0.05, y, w / len(vals) - 0.1, h, theme.surface)
            text_box(slide, canvas, labels[i] if i < len(labels) else "Step", bx + 0.12, y + 0.2, w / len(vals) - 0.24, 0.22, theme=theme, size=12, bold=True, fill=theme.accent, align=PP_ALIGN.CENTER)
            text_box(slide, canvas, value, bx + 0.14, y + 0.62, w / len(vals) - 0.28, h - 0.8, theme=theme, size=12, align=PP_ALIGN.CENTER)
