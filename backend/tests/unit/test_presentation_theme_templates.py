"""Coverage for theme packages, reusable templates, and patent components."""

from __future__ import annotations

from io import BytesIO

import pytest
from pptx import Presentation

from backend.app.presentation.contracts import PresentationDesign
from backend.app.presentation.pptx_renderer import render_pptx
from backend.app.presentation.renderer.theme import THEMES

pytestmark = pytest.mark.unit


THEME_IDS = [
    "patent_exam_classic",
    "legal_case_analysis",
    "technical_blueprint",
    "minimal_academic",
    "practice_workshop",
]


def _design(theme: str, template_id: str) -> PresentationDesign:
    return PresentationDesign.model_validate(
        {
            "title": "专利审查实务",
            "theme": theme,
            "slides": [
                {
                    "id": "slide_001",
                    "order": 1,
                    "layout": "legal_citation_focus",
                    "template_id": template_id,
                    "title": "规则与案例判断",
                    "body": "围绕权威课程材料识别法律规则。",
                    "bullets": ["保留法条依据", "区分事实与结论"],
                    "steps": ["申请", "审查", "授权"],
                    "legal_reference": "专利法第二十二条",
                    "legal_summary": "授予专利权的发明应当具备新颖性、创造性和实用性。",
                    "issue": "是否具备新颖性？",
                    "rule": "以权威规则为判断依据。",
                    "application": "逐项比对技术特征。",
                    "conclusion": "根据比对结果作出结论。",
                    "left_title": "权利要求",
                    "left_items": ["技术特征 A"],
                    "right_title": "现有技术",
                    "right_items": ["对应特征"],
                    "warning": "不要把初步审查等同于授权。",
                    "speaker_notes": "本页讲解规则和判断方法。",
                }
            ],
        }
    )


@pytest.mark.parametrize("theme", THEME_IDS)
def test_each_theme_package_renders(theme: str) -> None:
    assert theme in THEMES
    deck = Presentation(BytesIO(render_pptx(_design(theme, "legal_citation_focus"))))
    assert len(deck.slides) == 1
    assert len(deck.slides[0].shapes) >= 6
    assert deck.slides[0].notes_slide.notes_text_frame.text


@pytest.mark.parametrize(
    "template_id",
    ["cover_minimal", "content_rule_card", "irac_flow", "comparison_matrix", "timeline_process", "exam_checklist", "summary_roadmap"],
)
def test_patent_templates_render_native_objects(template_id: str) -> None:
    layout = "title" if template_id == "cover_minimal" else "content"
    design = _design("patent_exam_classic", template_id).model_copy(update={"slides": [
        _design("patent_exam_classic", template_id).slides[0].model_copy(update={"layout": layout})
    ]})
    slide = Presentation(BytesIO(render_pptx(design))).slides[0]
    assert len(slide.shapes) >= 4
    assert design.slides[0].speaker_notes in slide.notes_slide.notes_text_frame.text


def test_unknown_theme_is_rejected() -> None:
    with pytest.raises(ValueError):
        _design("unknown_theme", "content_rule_card")
