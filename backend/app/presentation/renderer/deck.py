"""Deck assembly using python-pptx native editable objects."""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.util import Inches

from backend.app.presentation.contracts import PresentationDesign
from backend.app.presentation.renderer.canvas import Canvas
from backend.app.presentation.renderer.layouts import render_slide
from backend.app.presentation.renderer.theme import theme_for


def render_design(design: PresentationDesign) -> bytes:
    presentation = Presentation()
    canvas = Canvas()
    presentation.slide_width = Inches(canvas.width)
    presentation.slide_height = Inches(canvas.height)
    blank = presentation.slide_layouts[6]
    theme = theme_for(design)
    for page, item in enumerate(design.slides, start=1):
        slide = presentation.slides.add_slide(blank)
        background = slide.background.fill
        background.solid()
        from pptx.dml.color import RGBColor

        background.fore_color.rgb = RGBColor.from_string(theme.background)
        render_slide(slide, canvas, item, theme, page)
        slide.notes_slide.notes_text_frame.text = item.speaker_notes
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
