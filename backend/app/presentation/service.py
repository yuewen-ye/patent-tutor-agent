"""LLM-designed, deterministically rendered complete PowerPoint artifacts."""

from __future__ import annotations

import hashlib
import json
import logging
import tempfile
from datetime import UTC, datetime
from io import BytesIO
from itertools import pairwise
from pathlib import Path
from typing import Any, Literal, cast, get_args
from zipfile import BadZipFile, ZipFile

from pptx import Presentation

from backend.app.agents.common import (
    generate_validated_json_stream,
    load_prompt,
)
from backend.app.core.agent_runtime_config import agent_temperature
from backend.app.core.llm import LLMClient, LLMMessage, LLMProviderError
from backend.app.presentation.contracts import (
    PresentationArtifact,
    PresentationCoursePackage,
    PresentationDesign,
    PresentationPreviewManifest,
    PresentationResult,
    PresentationSlide,
    PresentationSource,
    PresentationTemplate,
    PresentationVisualElement,
    PresentationVisualSlide,
    PresentationVisualStyle,
)
from backend.app.presentation.pptx_renderer import render_pptx
from backend.app.presentation.preview import generate_slide_previews

PPTX_MIME: Literal["application/vnd.openxmlformats-officedocument.presentationml.presentation"] = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
_LOGGER = logging.getLogger(__name__)
_PRESENTATION_SYSTEM = load_prompt(__file__)


def _strip_unknown_fields(raw: dict[str, Any], model: type[Any]) -> dict[str, Any]:
    """Drop fields that are not declared on ``model`` to survive ``extra='forbid'``."""
    allowed = set(model.model_fields.keys())
    return {k: v for k, v in raw.items() if k in allowed}


def _normalize_presentation_design(raw: object) -> object:
    """Normalize common LLM deviations before Pydantic validation.

    In addition to value-level fixes, we strip any extra fields the LLM added
    outside the ``PresentationDesign`` contract so that ``extra='forbid'`` does
    not degrade an otherwise valid deck.
    """
    if not isinstance(raw, dict):
        return raw

    design = _strip_unknown_fields(raw, PresentationDesign)

    if isinstance(design.get("visual_style"), dict):
        design["visual_style"] = _strip_unknown_fields(
            design["visual_style"], PresentationVisualStyle
        )

    slides = design.get("slides")
    if not isinstance(slides, list):
        return design

    allowed_slide = set(PresentationVisualSlide.model_fields.keys())
    allowed_element = set(PresentationVisualElement.model_fields.keys())
    normalized_slides: list[dict[str, Any]] = []
    for slide in slides:
        if not isinstance(slide, dict):
            normalized_slides.append(slide)
            continue
        normalized = {k: v for k, v in slide.items() if k in allowed_slide}

        legal_ref = normalized.get("legal_reference")
        if isinstance(legal_ref, list):
            parts = [str(item) for item in legal_ref if item]
            normalized["legal_reference"] = "; ".join(parts) if parts else None
        elif legal_ref == "":
            normalized["legal_reference"] = None

        composition = str(normalized.get("composition") or "")
        if composition == "timeline":
            normalized["composition"] = "timeline_with_callout"

        allowed_templates = set(get_args(PresentationTemplate))
        template_id = normalized.get("template_id")
        if template_id is not None and str(template_id) not in allowed_templates:
            normalized["template_id"] = None

        visual_elements = normalized.get("visual_elements")
        if isinstance(visual_elements, list):
            normalized_elements: list[dict[str, Any]] = []
            for element in visual_elements:
                if not isinstance(element, dict):
                    normalized_elements.append(element)
                    continue
                normalized_element = {
                    k: v for k, v in element.items() if k in allowed_element
                }
                element_type = str(normalized_element.get("type") or "")
                if element_type == "summary_roadmap":
                    normalized_element["type"] = "concept_map"
                normalized_elements.append(normalized_element)
            normalized["visual_elements"] = normalized_elements
        normalized_slides.append(normalized)
    design["slides"] = normalized_slides
    return design


def build_presentation_source(
    course_package: dict[str, Any], course_slides: dict[str, Any]
) -> PresentationSource:
    slides: list[PresentationSlide] = []
    mapping = course_slides.get("slide_to_block_id") or {}
    for item in course_slides.get("slides") or []:
        if not isinstance(item, dict):
            continue
        narration = item.get("narration") or {}
        slide_id = str(item.get("id") or f"slide_{len(slides) + 1:03d}")
        slides.append(
            PresentationSlide(
                id=slide_id,
                order=int(item.get("order") or len(slides) + 1),
                type=str(item.get("type") or "bullet"),
                title=str(item.get("title") or ""),
                content=cast(dict[str, object], item.get("content"))
                if isinstance(item.get("content"), dict)
                else {},
                narration=str(narration.get("text") or ""),
                source_block_id=str(mapping[slide_id]) if mapping.get(slide_id) else None,
            )
        )
    if not slides:
        raise ValueError("course_slides contains no slides")
    package = PresentationCoursePackage(
        title=str(course_package.get("title")) if course_package.get("title") else None,
        teaching_content=(
            str(course_package.get("teaching_content"))
            if course_package.get("teaching_content")
            else None
        ),
        legal_basis=(
            cast(list[object], course_package.get("legal_basis"))
            if isinstance(course_package.get("legal_basis"), list)
            else []
        ),
        block_plan=(course_package.get("block_plan") if isinstance(course_package.get("block_plan"), dict) else None),
        assessment=(course_package.get("assessment") if isinstance(course_package.get("assessment"), dict) else None),
    )
    return PresentationSource(course_package=package, slides=slides)


def _validate_design(design: PresentationDesign, source: PresentationSource) -> PresentationDesign:
    expected = [(slide.id, slide.order) for slide in source.slides]
    actual = [(slide.id, slide.order) for slide in design.slides]
    if actual != expected:
        raise ValueError("PresentationDesign must preserve every source slide id and order")
    templates = [slide.template_id or slide.layout for slide in design.slides]
    if len(design.slides) >= 4 and len(set(templates)) < 3:
        _LOGGER.warning("PresentationDesign uses fewer than three visual templates")
    if any(left == right for left, right in pairwise(templates)):
        _LOGGER.warning("Adjacent presentation slides reuse the same visual template")
    return design


def _validate_pptx(content: bytes, design: PresentationDesign) -> None:
    """Run package, Office parser, editable-shape and notes delivery checks.

    Package integrity and slide/note counts remain hard errors: a broken or
    incomplete PPTX is not usable. Content-level checks (title presence, shape
    count, exact speaker-notes match) are logged as warnings so that a single
    renderer imperfection does not degrade the whole artifact.
    """
    if not content or not content.startswith(b"PK"):
        raise ValueError("renderer did not return a PPTX zip package")
    try:
        with ZipFile(BytesIO(content)) as archive:
            names = set(archive.namelist())
            required = {"[Content_Types].xml", "ppt/presentation.xml", "ppt/theme/theme1.xml"}
            if not required.issubset(names):
                raise ValueError("renderer returned an incomplete PPTX package")
            notes = {name for name in names if name.startswith("ppt/notesSlides/notesSlide")}
            if len(notes) != len(design.slides):
                raise ValueError("renderer did not emit speaker notes for every slide")
    except BadZipFile as exc:
        raise ValueError("renderer returned an invalid PPTX zip package") from exc
    presentation = Presentation(BytesIO(content))
    if len(presentation.slides) != len(design.slides):
        raise ValueError("renderer output slide count differs from PresentationDesign")
    for source, rendered in zip(design.slides, presentation.slides, strict=True):
        text = " ".join(shape.text for shape in rendered.shapes if hasattr(shape, "text"))
        if source.title not in text:
            _LOGGER.warning("renderer may have lost title for slide %s", source.id)
        if not rendered.shapes:
            _LOGGER.warning("renderer emitted no editable shapes for slide %s", source.id)
        notes_text = rendered.notes_slide.notes_text_frame.text
        if not notes_text:
            _LOGGER.warning("renderer emitted empty speaker notes for slide %s", source.id)
        elif source.speaker_notes not in notes_text:
            _LOGGER.warning("renderer speaker notes differ from source for slide %s", source.id)


def generate_presentation_artifact(
    *,
    artifact_root: Path,
    session_id: str,
    course_package: dict[str, Any],
    course_slides: dict[str, Any],
    llm_client: LLMClient,
) -> dict[str, Any]:
    source = build_presentation_source(course_package, course_slides)
    messages = [
        LLMMessage(role="system", content=_PRESENTATION_SYSTEM),
        LLMMessage(
            role="user",
            content="请基于下列权威素材设计完整 PowerPoint：\n"
            + json.dumps(source.model_dump(), ensure_ascii=False, separators=(",", ":")),
        ),
    ]
    try:
        design = generate_validated_json_stream(
            llm_client,
            messages=messages,
            temperature=agent_temperature("generate_pptx", 0.2),
            agent="generate_pptx",
            output_model=PresentationDesign,
            schema_name="PresentationDesign",
            normalize=_normalize_presentation_design,
        )
    except LLMProviderError as exc:
        raise ValueError(f"Failed to generate presentation design: {exc}") from exc
    _validate_design(design, source)
    content = render_pptx(design)
    _validate_pptx(content, design)
    safe_session = "".join(ch if ch.isalnum() or ch in "-_" else "-" for ch in session_id)
    safe_session = safe_session.strip("-_") or "session"
    target_dir = artifact_root / "sessions" / safe_session / "presentation"
    target_dir.mkdir(parents=True, exist_ok=True)
    target = target_dir / "course_deck.pptx"
    with tempfile.NamedTemporaryFile(
        dir=target_dir, prefix=".course_deck-", suffix=".tmp", delete=False
    ) as tmp:
        tmp.write(content)
        temp_path = Path(tmp.name)
    temp_path.replace(target)

    preview_dir = target_dir / "previews"
    preview_result = generate_slide_previews(
        target, preview_dir, artifact_root=artifact_root
    )

    artifact = PresentationArtifact(
        artifact_id=f"{safe_session}-course-deck",
        path=f"artifacts/sessions/{safe_session}/presentation/course_deck.pptx",
        title="完整课程 PowerPoint",
        mime_type=PPTX_MIME,
        sha256=hashlib.sha256(content).hexdigest(),
        size_bytes=len(content),
        created_at=datetime.now(UTC).isoformat(),
    )
    (target_dir / "pptx_manifest.json").write_text(
        json.dumps(
            {
                "artifact": artifact.model_dump(),
                "source_slide_count": len(source.slides),
                "design_theme": design.theme,
                "preview_images": preview_result,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    return PresentationResult(
        status="generated",
        provider="configured_llm",
        source_slide_count=len(source.slides),
        speaker_notes_status="written",
        artifact=artifact,
        preview_images=PresentationPreviewManifest.model_validate(preview_result),
    ).model_dump()
