"""Native editable PPTX layout renderer coverage."""

from __future__ import annotations

from io import BytesIO
from zipfile import ZipFile

import pytest
from pptx import Presentation

from backend.app.presentation.contracts import PresentationDesign
from backend.app.presentation.pptx_renderer import render_pptx

pytestmark = pytest.mark.unit


def test_renderer_creates_editable_shapes_and_notes_for_layouts() -> None:
    layouts = ["title", "content", "two_column", "comparison", "process", "summary"]
    design = PresentationDesign.model_validate(
        {
            "title": "专利新颖性",
            "theme": "patent_blue",
            "slides": [
                {
                    "id": f"slide_{index:03d}",
                    "order": index,
                    "layout": layout,
                    "title": f"页面 {index}",
                    "subtitle": "副标题",
                    "body": "本页解释专利新颖性判断。",
                    "bullets": ["要点一", "要点二"],
                    "steps": ["检索", "比对", "结论"],
                    "left_title": "规则",
                    "left_items": ["单一现有技术", "完整公开"],
                    "right_title": "适用",
                    "right_items": ["确定时间", "逐项比对"],
                    "speaker_notes": f"第 {index} 页讲稿。",
                }
                for index, layout in enumerate(layouts, start=1)
            ],
        }
    )

    content = render_pptx(design)
    presentation = Presentation(BytesIO(content))

    assert len(presentation.slides) == len(layouts)
    assert all(len(slide.shapes) >= 3 for slide in presentation.slides)
    assert presentation.slides[2].shapes
    assert presentation.slides[4].notes_slide.notes_text_frame.text == "第 5 页讲稿。"
    with ZipFile(BytesIO(content)) as archive:
        assert len([n for n in archive.namelist() if n.startswith("ppt/notesSlides/notesSlide")]) == 6
